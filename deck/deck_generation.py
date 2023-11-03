import os
import openai
from pptx import Presentation
from pptx.util import Inches
import random 
import re
from langchain.prompts import PromptTemplate
# from elysium_prompts.deck_prompts.pitch_deck_prompts import personalized_pitch_deck_prompt
import requests
from elysium_prompts.deck_prompts.pitch_deck_prompts import personalized_pitch_deck_prompt, presentation_prompt
import os, logging
# openai.api_key = os.getenv("OPENAI_API_KEY")
openai.api_key = "YOUR_OPENAI_KEY"


SLIDE_WIDTH = Inches(10)
SLIDE_HEIGHT = Inches(5.625)

# TEXT_LEFT = Inches(1)  
# TEXT_WIDTH = Inches(5)
# # TEXT_HEIGHT = Inches(3)

# # IMAGE_RIGHT = Inches(6.5)
# # IMAGE_WIDTH = Inches(4) 
# # IMAGE_HEIGHT = Inches(3)

# TEXT_TOP = Inches(0.5)
# TEXT_LEFT = Inches(1.5)
# TEXT_WIDTH = Inches(6)
# TEXT_HEIGHT = Inches(2.5)

# IMAGE_TOP = Inches(3)
# IMAGE_LEFT = Inches(0.5)
# IMAGE_WIDTH = Inches(8)
# IMAGE_HEIGHT = Inches(2.5)

def generate_image(text):
    response = openai.Image.create(
        n=1,
        prompt="futurisitc AI robot, trending on art station",
        size="256x256"
    )
    image_url = response['data'][0]['url']
    image_data = requests.get(image_url).content
    return image_data

def create_ppt_text(topic, recipient_name, entity_name, entity_info="", recent_tweets="", custom_prompt = None):
    if custom_prompt:
        prompt = str(presentation_prompt)
    else:
        prompt = PromptTemplate.from_template(personalized_pitch_deck_prompt)
        formatted_prompt = prompt.format(
            topic="", 
            recipient_name=recipient_name, 
            entity_name=entity_name, 
            entity_info=entity_info,
            recent_tweets=recent_tweets
        )
    if custom_prompt:
        response = openai.ChatCompletion.create(
            model="gpt-3.5-turbo-16k",
            messages=[
                {"role": "system", "content": (prompt)},
                {"role": "user", "content": ("The user wants a presentation about, TOPIC:" + topic + ", INSTRUCTIONS: " + custom_prompt )},
            ],
            temperature=0.5,
        )
    else:
        response = openai.ChatCompletion.create(
            model="gpt-3.5-turbo-16k",
            messages=[
                {"role": "system", "content": (formatted_prompt)},
                {"role": "user", "content": ("The user wants a presentation about " + topic )}
            ],
            temperature=0.5,
        )

    return response['choices'][0]['message']['content']

# def create_ppt_v1(text_file, design_number, ppt_name):
#     prs = Presentation(f"deck/Designs/Design-{design_number}.pptx")
#     slide_count = 0
#     header = ""
#     content = ""
#     last_slide_layout_index = -1
#     firsttime = True
#     text_on_left = True

#     with open(text_file, 'r', encoding='utf-8') as f:
#         for line_num, line in enumerate(f):
#             if line.startswith('#Title:'):
#                 header = line.replace('#Title:', '').strip()
#                 slide = prs.slides.add_slide(prs.slide_layouts[0])
#                 title = slide.shapes.title
#                 title.text = header
#                 body_shape = slide.shapes.placeholders[1]
#                 continue
#             elif line.startswith('#Slide:'):
#                 if slide_count > 0:
#                     slide = prs.slides.add_slide(prs.slide_layouts[slide_layout_index])
#                     title = slide.shapes.title
#                     title.text = header
#                     body_shape = slide.shapes.placeholders[slide_placeholder_index]
#                     tf = body_shape.text_frame
#                     tf.text = content
#                 content = "" 
#                 slide_count += 1
#                 slide_layout_index = last_slide_layout_index
#                 layout_indices = [1, 7, 8] 
#                 while slide_layout_index == last_slide_layout_index:
#                     if firsttime == True:
#                         slide_layout_index = 1
#                         slide_placeholder_index = 1
#                         firsttime = False
#                         break
#                     slide_layout_index = random.choice(layout_indices) # Select random slide index
#                     if slide_layout_index == 8:
#                         slide_placeholder_index = 2
#                     else:
#                         slide_placeholder_index = 1
#                 last_slide_layout_index = slide_layout_index
#                 continue

#             elif line.startswith('#Header:'):
#                 header = line.replace('#Header:', '').strip()
#                 continue

#             elif line.startswith('#Content:'):
#                 content = line.replace('#Content:', '').strip()
#                 next_line = f.readline().strip()
#                 while next_line and not next_line.startswith('#'):
#                     content += '\n' + next_line
#                     next_line = f.readline().strip()

#                 image_data = generate_image(content)
#                 image_path = f'deck/Images/slide_{slide_count}.png'
#                 with open(image_path, 'wb') as image_file:
#                     image_file.write(image_data)

#                 # Determine placement based on alternating pattern
#                 if text_on_left:
#                     text_left = Inches(0.5)
#                     image_left = SLIDE_WIDTH - IMAGE_WIDTH - Inches(0.5)
#                 else:
#                     image_left = Inches(0.5)
#                     text_left = SLIDE_WIDTH - TEXT_WIDTH - Inches(0.5)
                
#                 # Add text content
#                 text_box = slide.shapes.add_textbox(text_left, Inches(0.5), TEXT_WIDTH, TEXT_HEIGHT)
#                 tf = text_box.text_frame
#                 tf.text = content

#                 # Add image
#                 slide.shapes.add_picture(image_path, image_left, Inches(0.5))

#                 # Alternate side for next slide
#                 text_on_left = not text_on_left
#                 # Add image to the existing slide
#                 # slide.shapes.add_picture(image_path, image_left, image_top, width=Inches(6), height=Inches(4))
#                 # Add image to slide
#                 # slide = prs.slides.add_slide(prs.slide_layouts[slide_layout_index])
#                 # slide.shapes.add_picture(image_path, Inches(1), Inches(1), width=Inches(6), height=Inches(4))
#                 continue
                

#     prs.save(f'deck/GeneratedPresentations/{ppt_name}.pptx')
#     file_path = f"GeneratedPresentations/{ppt_name}.pptx"

# def create_ppt_v2(content, design_number, ppt_name):
#     lines = [line.strip() for line in content.split("\n") if line.strip()]
#     prs = Presentation(f"deck/Designs/Design-{design_number}.pptx")
#     slide_content = ""
#     title_text = ""

#     # Image positioning
#     image_left_flag = True  # Start with image on the left
#     img_left_position = Inches(0.5)
#     img_right_position = Inches(8)
#     img_top_position = Inches(2)
#     img_width = Inches(5)
    
#     for idx, line in enumerate(lines):
#         if line.endswith(":::") and not slide_content:
#             title_text = line.replace(":::", "")
#         elif line.endswith(":::") and slide_content:
#             slide = prs.slides.add_slide(prs.slide_layouts[1])
#             title = slide.shapes.title
#             subtitle = slide.placeholders[1]
#             title.text = title_text
#             subtitle.text_frame.text = slide_content.strip()  # changed to text_frame for indentation
#             p = subtitle.text_frame.add_paragraph()  # added for indentation
#             p.text = slide_content.strip()
#             p.level = 1 
#             # Add image to the slide
#             image_data = generate_image(slide_content)
#             image_path = f'deck/Images/slide_{idx}.png'
#             with open(image_path, 'wb') as image_file:
#                     image_file.write(image_data)
#             if image_left_flag:
#                 slide.shapes.add_picture(image_path, img_left_position, img_top_position, img_width)
#             else:
#                 slide.shapes.add_picture(image_path, img_right_position, img_top_position, img_width)
            
#             # Toggle the image position for the next slide
#             image_left_flag = not image_left_flag

#             title_text = line.replace(":::", "")
#             slide_content = ""
#         else:
#             slide_content += line + "\n"

#     # For any remaining content
#     if slide_content:
#         slide = prs.slides.add_slide(prs.slide_layouts[1])
#         title = slide.shapes.title
#         subtitle = slide.placeholders[1]
#         title.text = title_text
#         subtitle.text_frame.text = slide_content.strip()
#         p = subtitle.text_frame.add_paragraph()
#         p.text = slide_content.strip()
#         p.level = 1

#         # Add image to the slide
#         image_data = generate_image(slide_content)
#         image_path = f'deck/Images/slide_{idx}.png'
#         with open(image_path, 'wb') as image_file:
#                 image_file.write(image_data)
#         if image_left_flag:
#             slide.shapes.add_picture(image_path, img_left_position, img_top_position, img_width)
#         else:
#             slide.shapes.add_picture(image_path, img_right_position, img_top_position, img_width)

#     prs.save(f'deck/GeneratedPresentations/{ppt_name}.pptx')

# # NOTE: Best so far
# def create_ppt_v3(text_file, design_number, ppt_name):
    prs = Presentation(f"deck/Designs/Design-{design_number}.pptx")
    slide_count = 0
    header = ""
    content = ""
    last_slide_layout_index = -1
    firsttime = True

    # Positioning constants
    TEXT_LEFT = Inches(0.5)
    TEXT_RIGHT = Inches(5.5)
    TEXT_WIDTH = Inches(5)
    TEXT_HEIGHT = Inches(5)

    IMAGE_LEFT = Inches(0.5)
    IMAGE_RIGHT = Inches(5.5)
    IMAGE_WIDTH = Inches(5)
    IMAGE_HEIGHT = Inches(5)
    
    text_left_flag = True  # Start with text on the left

    with open(text_file, 'r', encoding='utf-8') as f:
        for line_num, line in enumerate(f):
            if line.startswith('#Title:'):
                header = line.replace('#Title:', '').strip()
                slide = prs.slides.add_slide(prs.slide_layouts[0])
                title = slide.shapes.title
                title.text = header
                body_shape = slide.shapes.placeholders[1]
                continue
            elif line.startswith('#Slide:'):
                if slide_count > 0:
                    slide = prs.slides.add_slide(prs.slide_layouts[slide_layout_index])
                    title = slide.shapes.title
                    title.text = header
                    
                    # Adjust position based on the flag
                    if text_left_flag:
                        text_left_position = TEXT_LEFT
                        image_left_position = IMAGE_RIGHT
                    else:
                        text_left_position = TEXT_RIGHT
                        image_left_position = IMAGE_LEFT

                    # Add text to the slide
                    txBox = slide.shapes.add_textbox(text_left_position, Inches(1), TEXT_WIDTH, TEXT_HEIGHT)
                    tf = txBox.text_frame
                    tf.text = content

                    # Add image to the slide
                    image_data = generate_image(content)  # Assuming you've already defined generate_image function
                    image_path = f'deck/Images/slide_{slide_count}.png'
                    with open(image_path, 'wb') as image_file:
                        image_file.write(image_data)
                    slide.shapes.add_picture(image_path, image_left_position, Inches(1), IMAGE_WIDTH, IMAGE_HEIGHT)
                    
                    # Toggle the text position for the next slide
                    text_left_flag = not text_left_flag

                content = "" 
                slide_count += 1
                slide_layout_index = last_slide_layout_index
                layout_indices = [1, 7, 8] 
                while slide_layout_index == last_slide_layout_index:
                    if firsttime == True:
                        slide_layout_index = 1
                        slide_placeholder_index = 1
                        firsttime = False
                        break
                    slide_layout_index = random.choice(layout_indices)
                    if slide_layout_index == 8:
                        slide_placeholder_index = 2
                    else:
                        slide_placeholder_index = 1
                last_slide_layout_index = slide_layout_index
                continue
            elif line.startswith('#Header:'):
                header = line.replace('#Header:', '').strip()
                continue
            elif line.startswith('#Content:'):
                content = line.replace('#Content:', '').strip()
                next_line = f.readline().strip()
                while next_line and not next_line.startswith('#'):
                    content += '\n' + next_line
                    next_line = f.readline().strip()
                continue
    prs.save(f'deck/GeneratedPresentations/{ppt_name}.pptx')
    file_path = f"deck/GeneratedPresentations/{ppt_name}.pptx"


# def create_ppt_v4(
#         text_file, 
#         design_number, 
#         ppt_name,
#         twitter_pfp_img_url = None
# ):
#     prs = Presentation(f"deck/Designs/Design-{design_number}.pptx")
#     slide_count = 0
#     header = ""
#     content = ""

#     # Positioning constants
#     TEXT_LEFT = Inches(0.5)
#     TEXT_WIDTH = Inches(5)
#     TEXT_HEIGHT = Inches(5)
    
#     IMAGE_BOTTOM = Inches(3.5)
#     IMAGE_RIGHT = Inches(8.5)
#     IMAGE_WIDTH = Inches(4)
#     IMAGE_HEIGHT = Inches(3)

#     with open(text_file, 'r', encoding='utf-8') as f:
#         for line_num, line in enumerate(f):
#             if line.startswith('#Title:'):
#                 header = line.replace('#Title:', '').strip()
#                 slide = prs.slides.add_slide(prs.slide_layouts[0])
#                 title = slide.shapes.title
#                 title.text = header
#                 continue
#             elif line.startswith('#Slide:'):
#                 if slide_count > 0:
#                     slide = prs.slides.add_slide(prs.slide_layouts[1])
#                     title = slide.shapes.title
#                     title.text = header

#                     # Add text to the slide
#                     txBox = slide.shapes.add_textbox(TEXT_LEFT, Inches(1.5), TEXT_WIDTH, TEXT_HEIGHT)
#                     tf = txBox.text_frame
#                     tf.text = content

#                     # Add image to the slide
#                     image_data = generate_image(content)  # Assuming you've already defined generate_image function
#                     image_path = f'deck/Images/slide_{slide_count}.png'
#                     with open(image_path, 'wb') as image_file:
#                         image_file.write(image_data)
#                     slide.shapes.add_picture(image_path, IMAGE_RIGHT, IMAGE_BOTTOM, IMAGE_WIDTH, IMAGE_HEIGHT)

#                 content = "" 
#                 slide_count += 1
#                 continue
#             elif line.startswith('#Header:'):
#                 header = line.replace('#Header:', '').strip()
#                 continue
#             elif line.startswith('#Content:'):
#                 content = line.replace('#Content:', '').strip()
#                 next_line = f.readline().strip()
#                 while next_line and not next_line.startswith('#'):
#                     content += '\n' + next_line
#                     next_line = f.readline().strip()
#                 continue
#     prs.save(f'deck/GeneratedPresentations/{ppt_name}.pptx')
#     file_path = f"deck/GeneratedPresentations/{ppt_name}.pptx"
#     return file_path

IMAGE_FOLDER = 'deck/Images' 

def get_random_image():
  images = os.listdir(IMAGE_FOLDER)
  image = random.choice(images)
  return os.path.join(IMAGE_FOLDER, image)

from io import BytesIO

def create_ppt_v5(
        text_file, 
        design_number, 
        ppt_name,
        twitter_pfp_img_url = None,
):
    prs = Presentation(f"deck/Designs/Design-{design_number}.pptx")
    slide_count = 0
    header = ""
    content = ""

    # Positioning constants
    TEXT_LEFT = Inches(0.5)
    TEXT_TOP = Inches(3.5) # Moved down 1 inch
    TEXT_WIDTH = Inches(5) 
    TEXT_HEIGHT = Inches(5)

    # Later when adding text box:
    
    IMAGE_BOTTOM = Inches(3.5)
    IMAGE_RIGHT = Inches(8.5)
    IMAGE_WIDTH = Inches(4)
    IMAGE_HEIGHT = Inches(3)

    with open(text_file, 'r', encoding='utf-8') as f:
        for line_num, line in enumerate(f):
            if line.startswith('#Title:'):
                title_content = line.replace('#Title:', '').strip()
                slide = prs.slides.add_slide(prs.slide_layouts[0])  # Assuming layout 0 is for the title
                title = slide.shapes.title
                title.text = title_content
                continue
            if line.startswith('#Slide:'):
                slide_count += 1
                continue

            elif line.startswith('#Header:'):
                header = line.replace('#Header:', '').strip()
                continue

            elif line.startswith('#Content:'):
                content = line.replace('#Content:', '').strip()
                next_line = f.readline().strip()
                while next_line and not next_line.startswith('#'):
                    content += '\n' + next_line
                    next_line = f.readline().strip()

                slide = prs.slides.add_slide(prs.slide_layouts[1])
                title = slide.shapes.title
                title.text = header

                # Add text to the slide
                txBox = slide.shapes.add_textbox(TEXT_LEFT, TEXT_TOP, TEXT_WIDTH, TEXT_HEIGHT)
                tf = txBox.text_frame
                tf.text = content
               

                # Add image to the slide NOTE: commented out to save time and tokens right now - merk 9/12/23
                # image_data = generate_image(content)
                # image_path = f'deck/Images/slide_{slide_count}.png'
                # with open(image_path, 'wb') as image_file:
                #     image_file.write(image_data)
                image_path = get_random_image()

                left = random.randint(1, 7)
                top = random.randint(1, 5)
                # slide.shapes.add_picture(image_path, Inches(left), Inches(top))
                
                slide.shapes.add_picture(image_path, IMAGE_RIGHT, IMAGE_BOTTOM, IMAGE_WIDTH, IMAGE_HEIGHT)

                # If it's slide 2 and there's a twitter profile image URL, add the image
                if slide_count == 2 and twitter_pfp_img_url:
                    response = requests.get(twitter_pfp_img_url)
                    image_stream = BytesIO(response.content)
                    slide.shapes.add_picture(image_stream, IMAGE_RIGHT, IMAGE_BOTTOM, IMAGE_WIDTH, IMAGE_HEIGHT)

    current_directory = os.getcwd()
    logging.info(f"Current directory: {current_directory}")
    file_path = os.path.join(current_directory, "presentation.pptx")

    prs.save(file_path)
    return file_path
